import re
import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pyodbc
import win32com.client as win32

def extraer_ips(texto):
    patron_ipv4 = r'\b(?:\d{1,3}\.){3}\d{1,3}\b'
    patron_ipv6 = r'\b(?:[0-9a-fA-F]{1,4}:){7}[0-9a-fA-F]{1,4}\b'
    patron = f'({patron_ipv4})|({patron_ipv6})'
    return [ip[0] or ip[1] for ip in re.findall(patron, texto)]

def extraer_ips_txt(ruta_archivo):
    try:
        with open(ruta_archivo, 'r', encoding='utf-8') as archivo:
            texto = archivo.read()
        return extraer_ips(texto)
    except FileNotFoundError as e:
        print(f"Error al leer el archivo TXT: {e}")
        return []

def extraer_ips_pdf(ruta_archivo):
    try:
        with open(ruta_archivo, 'rb') as archivo:
            lector = PdfReader(archivo)
            texto = ''.join(pagina.extract_text() or '' for pagina in lector.pages)
        return extraer_ips(texto)
    except Exception as e:
        print(f"Error al leer el archivo PDF: {e}")
        return []

def extraer_ips_docx(ruta_archivo):
    try:
        documento = Document(ruta_archivo)
        texto = '\n'.join(parrafo.text for parrafo in documento.paragraphs)
        return extraer_ips(texto)
    except Exception as e:
        print(f"Error al leer el archivo DOCX: {e}")
        return []

def extraer_ips_excel(ruta_archivo):
    try:
        xls = pd.ExcelFile(ruta_archivo)
        ips = []
        for nombre_hoja in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=nombre_hoja)
            texto = df.to_string()
            ips.extend(extraer_ips(texto))
        return ips
    except Exception as e:
        print(f"Error al leer el archivo Excel: {e}")
        return []

def extraer_ips_pptx(ruta_archivo):
    try:
        presentacion = Presentation(ruta_archivo)
        texto = ''
        for diapositiva in presentacion.slides:
            for forma in diapositiva.shapes:
                if hasattr(forma, "text"):
                    texto += forma.text + '\n'
        return extraer_ips(texto)
    except Exception as e:
        print(f"Error al leer el archivo PPTX: {e}")
        return []

def extraer_ips_access(ruta_archivo):
    try:
        ips = []
        conn_str = r'DRIVER={Microsoft Access Driver (*.mdb, *.accdb)};DBQ=' + ruta_archivo
        conn = pyodbc.connect(conn_str)
        cursor = conn.cursor()
        cursor.tables()
        for table_info in cursor.tables(tableType='TABLE'):
            table_name = table_info.table_name
            df = pd.read_sql(f'SELECT * FROM {table_name}', conn)
            texto = df.to_string()
            ips.extend(extraer_ips(texto))
        conn.close()
        return ips
    except Exception as e:
        print(f"Error al leer el archivo Access: {e}")
        return []

def extraer_ips_publisher(ruta_archivo):
    try:
        pub = win32.Dispatch('Publisher.Application')
        doc = pub.Open(ruta_archivo)
        texto = doc.Pages(1).Shapes(1).TextFrame.TextRange.Text
        ips = extraer_ips(texto)
        doc.Close()
        pub.Quit()
        return ips
    except Exception as e:
        print(f"Error al leer el archivo Publisher: {e}")
        return []

def extraer_ips_onenote(ruta_archivo):
    try:
        onenote = win32.Dispatch("OneNote.Application")
        ns = r"namespaces.msoverlay://<namespace>"
        one_ns = onenote.GetNamespace(ns)
        onenote.OpenHierarchy(ruta_archivo, one_ns)
        texto = onenote.GetPageContent(ruta_archivo)
        ips = extraer_ips(texto)
        return ips
    except Exception as e:
        print(f"Error al leer el archivo OneNote: {e}")
        return []

def extraer_ips_de_archivo(ruta_archivo):
    extension = os.path.splitext(ruta_archivo)[1].lower()
    if extension == '.txt':
        return extraer_ips_txt(ruta_archivo)
    elif extension == '.pdf':
        return extraer_ips_pdf(ruta_archivo)
    elif extension == '.docx':
        return extraer_ips_docx(ruta_archivo)
    elif extension in ['.xls', '.xlsx']:
        return extraer_ips_excel(ruta_archivo)
    elif extension == '.pptx':
        return extraer_ips_pptx(ruta_archivo)
    elif extension in ['.mdb', '.accdb']:
        return extraer_ips_access(ruta_archivo)
    elif extension == '.pub':
        return extraer_ips_publisher(ruta_archivo)
    elif extension == '.one':
        return extraer_ips_onenote(ruta_archivo)
    else:
        print("Formato de archivo no soportado.")
        return []

# Ejemplo de uso de IP
ruta_archivo = "datos.txt"
ips_encontradas = extraer_ips_de_archivo(ruta_archivo)

# Aquí se imprimen las IP encontradas
print("Direcciones IP encontradas:")
for ip in ips_encontradas:
    print(ip)

